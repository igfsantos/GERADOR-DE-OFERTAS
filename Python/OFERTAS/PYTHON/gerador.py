from pptx import Presentation
from pptx.util import Inches, Pt
'''import mysql.connector'''

#banco de dados


'''conexao = mysql.connector(
    host='localhost',
    user='root',
    password='',
    database='bdofertas'
)
cursor = conexao.cursor()

#CRUD

cursor.close()
conexao.close()'''

#IMAGENS DO PROJETO
img_path = 'base.png' #IMAGEM DA BOA DO DIA
img_path2 = 'coringa.png' # IMAGEM DO PRODUTO QUE VAI SER DO BANCO DE DADOS

#Criando PowerPoint

apresentacao = Presentation()

#Criando Slide
slide1 = apresentacao.slides.add_slide(apresentacao.slide_layouts[0])

#colocando Imagem boa do dia
left = top = Inches(0)
pic = slide1.shapes.add_picture(img_path,
                               left, top)
left = Inches(5)
height = Inches(4)

#COLOCANDO A DATA
x = Inches(3.6) #Horizontal
y = Inches(1.9) #vertical
largura = Inches(2)
altura = Inches(2)


caixa_texto = slide1.shapes.add_textbox(x, y, largura, altura)
text_frame = caixa_texto.text_frame
data = text_frame.add_paragraph()
data.text = str(input('Qual data da oferta?  '))
data.name = 'Calibri'
data.font.bold = True
data.font.size = Pt(30)

#CAIXA DE TEXTO DO NOME DA MERCADORIA
x = Inches(4)
y = Inches(2.5)
largura = Inches(2)
altura = Inches(2)

caixa_texto1 = slide1.shapes.add_textbox(x, y, largura, altura)
text_frame = caixa_texto.text_frame
codigo = text_frame.add_paragraph()
codigo.text = str(input('Qual codigo interno  '))
codigo.name = 'Calibri'
codigo.font.bold = True
codigo.font.size = Pt(60)



#CAIXA DE TEXTO DO VALOR 0-9 REAIS
x = Inches(6.2) #Horizontal
y = Inches(4) #vertical
largura = Inches(2)
altura = Inches(2)

caixa_texto2 = slide1.shapes.add_textbox(x, y, largura, altura)
text_frame = caixa_texto2.text_frame
valor = text_frame.add_paragraph()
valor.text = str(input('Qual Valor? R$  '))
valor.font.bold = True
valor.font.size = Pt(120)

#CAIXA DE TEXTO DO R$ 0-9 REAIS
x = Inches(5.5) #Horizontal
y = Inches(5.4) #vertical
largura = Inches(2)
altura = Inches(2)

caixa_texto3 = slide1.shapes.add_textbox(x, y, largura, altura)
text_frame = caixa_texto3.text_frame
rs = text_frame.add_paragraph()
rs.text = ( 'R$')
rs.font.bold = True
rs.font.size = Pt(30)



#caixa de texto do 10--99



#caixa de texto de 100+






#colocando Imagem bo Encarte
x = Inches(1) #Horizontal
y = Inches(3.5) #vertical
largura = Inches(2)
altura = Inches(2)


pic = slide1.shapes.add_picture(img_path2,
                                x, y, largura, altura )



print('Sua arte foi Feita e enviada para o G-mail.')






apresentacao.save('OFERTAS.pptx')